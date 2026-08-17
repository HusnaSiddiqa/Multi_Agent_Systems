"""PPT Agent - Data Builders & PowerPoint Renderer.

Data builders reshape DataFrames into chart-ready structures.
The renderer creates native editable PowerPoint charts using python-pptx.
"""

import io
import os
import time
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LEGEND_POSITION

from .registry import CHART_REGISTRY


# ─── Data Builders ───────────────────────────────────────────────────────────

def _apply_category_transform(series: pd.Series, transform: str) -> pd.Series:
    """Apply date bucketing transform to a category series."""
    if transform == "none" or not transform:
        return series.astype(str)
    try:
        dt = pd.to_datetime(series, errors="coerce")
        if dt.isna().all():
            return series.astype(str)
        if transform == "quarter":
            return dt.dt.to_period("Q").astype(str)
        elif transform == "month":
            return dt.dt.to_period("M").astype(str)
        elif transform == "year":
            return dt.dt.year.astype(str)
    except Exception:
        pass
    return series.astype(str)


def _to_numeric_series(series: pd.Series) -> pd.Series:
    """Coerce a series to numeric, stripping formatting chars."""
    if pd.api.types.is_numeric_dtype(series):
        return series.fillna(0)
    cleaned = series.astype(str).str.replace(r"[\$,%]", "", regex=True)
    return pd.to_numeric(cleaned, errors="coerce").fillna(0)




def _sort_categories_asc(temp: pd.DataFrame, df: pd.DataFrame, enc: dict, cat_col: str) -> pd.DataFrame:
    """Sort a temp DataFrame by its 'cat' column in chronological ascending order.
    
    Detects quarter patterns (Q1'24), date strings, and explicit sort columns.
    Returns the sorted DataFrame (or unchanged if no date pattern detected).
    """
    import re as _re

    _sort_field = enc.get("category_sort_column", "auto")
    if not _sort_field or _sort_field.lower() in ("none", ""):
        return temp

    _did_sort = False

    # Try explicit sort column from LLM
    if _sort_field.lower() != "auto":
        _mcols = [c for c in df.columns if c.lower() == _sort_field.lower()]
        if _mcols and _mcols[0].lower() != cat_col.lower():
            _smap = df.drop_duplicates(subset=[cat_col]).set_index(cat_col)[_mcols[0]]
            temp["_sk"] = temp["cat"].map(_smap)
            temp = temp.sort_values("_sk").drop(columns=["_sk"])
            _did_sort = True

    if not _did_sort:
        # Auto: try quarter patterns (Q1'24, Q4 2025, 2024-Q1)
        _cats = temp["cat"].drop_duplicates().tolist()
        _qre = _re.compile(r"Q([1-4]).*?(\d{2,4})", _re.IGNORECASE)
        if any(_qre.search(str(c)) for c in _cats):
            _km = {}
            for c in _cats:
                m = _qre.search(str(c))
                if m:
                    q, y = int(m.group(1)), int(m.group(2))
                    _km[c] = (y + 2000 if y < 100 else y) * 10 + q
                else:
                    _km[c] = 0
            temp["_sk"] = temp["cat"].map(_km)
            temp = temp.sort_values("_sk").drop(columns=["_sk"])
            _did_sort = True

        if not _did_sort:
            # Try standard date parsing
            try:
                _p = pd.to_datetime(pd.Series(_cats), format="mixed", dayfirst=False)
                if not _p.isna().all():
                    temp["_sk"] = temp["cat"].map(dict(zip(_cats, _p)))
                    temp = temp.sort_values("_sk").drop(columns=["_sk"])
            except Exception:
                pass

    return temp




def build_category(df: pd.DataFrame, enc: dict) -> CategoryChartData:
    """Build single-series category chart data (pie, doughnut, simple bar/column).
    
    Channels: category, value
    """
    cat_col = enc["channels"]["category"]
    val_col = enc["channels"]["value"]
    transform = enc.get("category_transform", "none")
    aggregate = enc.get("aggregate", "sum")

    categories = _apply_category_transform(df[cat_col], transform)
    values = _to_numeric_series(df[val_col])

    # Aggregate if duplicates
    temp = pd.DataFrame({"cat": categories, "val": values})
    if aggregate == "mean":
        grouped = temp.groupby("cat", sort=False)["val"].mean()
    elif aggregate == "count":
        grouped = temp.groupby("cat", sort=False)["val"].count()
    else:
        grouped = temp.groupby("cat", sort=False)["val"].sum()

    chart_data = CategoryChartData()
    chart_data.categories = [str(c) for c in grouped.index]
    chart_data.add_series("Value", [float(v) for v in grouped.values])
    return chart_data


def build_pivot(df: pd.DataFrame, enc: dict) -> CategoryChartData:
    """Build multi-series pivoted chart data (column, bar, line, area, stacked).
    
    Channels: category, value, optional series
    """
    cat_col = enc["channels"]["category"]
    val_col = enc["channels"]["value"]
    series_col = enc["channels"].get("series")
    transform = enc.get("category_transform", "none")
    aggregate = enc.get("aggregate", "sum")

    categories = _apply_category_transform(df[cat_col], transform)
    values = _to_numeric_series(df[val_col])

    if not series_col or series_col not in df.columns:
        # Single series — same as build_category but returns CategoryChartData
        temp = pd.DataFrame({"cat": categories, "val": values})
        agg_func = {"mean": "mean", "count": "count"}.get(aggregate, "sum")
        grouped = temp.groupby("cat", sort=False)["val"].agg(agg_func)

        # Sort categories chronologically (ASC)
        sorted_idx = _sort_categories_asc(
            pd.DataFrame({"cat": grouped.index}), df, enc, cat_col
        )["cat"].tolist()
        grouped = grouped.reindex(sorted_idx)

        chart_data = CategoryChartData()
        chart_data.categories = [str(c) for c in grouped.index]
        chart_data.add_series(val_col, [float(v) for v in grouped.values])
        return chart_data

    # Multi-series pivot
    temp = pd.DataFrame({
        "cat": categories,
        "series": df[series_col].astype(str),
        "val": values,
    })

    # ── Sort categories chronologically (ASC) for chart axis ──
    temp = _sort_categories_asc(temp, df, enc, cat_col)

    agg_func = {"mean": "mean", "count": "count"}.get(aggregate, "sum")
    pivot = temp.pivot_table(index="cat", columns="series", values="val",
                             aggfunc=agg_func, fill_value=0, sort=False)

    chart_data = CategoryChartData()
    chart_data.categories = [str(c) for c in pivot.index]
    for col in pivot.columns:
        chart_data.add_series(str(col), [float(v) for v in pivot[col].values])
    return chart_data


def build_xy(df: pd.DataFrame, enc: dict) -> XyChartData:
    """Build XY scatter chart data.
    
    Channels: x, y
    """
    x_col = enc["channels"]["x"]
    y_col = enc["channels"]["y"]

    x_vals = _to_numeric_series(df[x_col])
    y_vals = _to_numeric_series(df[y_col])

    chart_data = XyChartData()
    series = chart_data.add_series("Data")
    for x, y in zip(x_vals, y_vals):
        series.add_data_point(float(x), float(y))
    return chart_data


def build_bubble(df: pd.DataFrame, enc: dict) -> BubbleChartData:
    """Build bubble chart data.
    
    Channels: x, y, size
    """
    x_col = enc["channels"]["x"]
    y_col = enc["channels"]["y"]
    size_col = enc["channels"]["size"]

    x_vals = _to_numeric_series(df[x_col])
    y_vals = _to_numeric_series(df[y_col])
    size_vals = _to_numeric_series(df[size_col])

    chart_data = BubbleChartData()
    series = chart_data.add_series("Data")
    for x, y, s in zip(x_vals, y_vals, size_vals):
        series.add_data_point(float(x), float(y), float(s))
    return chart_data


# Map builder name -> function
DATA_BUILDERS = {
    "build_category": build_category,
    "build_pivot": build_pivot,
    "build_xy": build_xy,
    "build_bubble": build_bubble,
}


# ─── Renderer ────────────────────────────────────────────────────────────────

# ─── Professional color palette ──────────────────────────────────────────────

# Client design guide — text/layout colors only
NAVY       = RGBColor(0x20, 0x36, 0x61)

# Standard chart colors (medium tone — universally readable, Excel/Tableau style)
CHART_COLORS = [
    RGBColor(0x4E, 0x79, 0xA7),  # Steel Blue
    RGBColor(0xF2, 0x8E, 0x2B),  # Orange
    RGBColor(0xE1, 0x57, 0x59),  # Coral Red
    RGBColor(0x76, 0xB7, 0xB2),  # Teal
    RGBColor(0x59, 0xA1, 0x4F),  # Green
    RGBColor(0xED, 0xC9, 0x48),  # Gold
    RGBColor(0xAF, 0x7A, 0xA1),  # Purple
    RGBColor(0xFF, 0x9D, 0xA7),  # Pink
    RGBColor(0x9C, 0x75, 0x5F),  # Brown
    RGBColor(0xBA, 0xB0, 0xAC),  # Warm Gray
]

DARK_GRAY  = RGBColor(0x54, 0x56, 0x5B)   # body text
MID_GRAY   = RGBColor(0x96, 0x99, 0x9F)   # captions/footer
FONT_NAME  = "Trebuchet MS"


def _is_percentage_data(df: pd.DataFrame, enc: dict) -> bool:
    """Detect if the value column contains percentage/share data (0-100 range)."""
    val_col = enc.get("channels", {}).get("value")
    if not val_col or val_col not in df.columns:
        return False
    numeric = pd.to_numeric(df[val_col], errors="coerce").dropna()
    if numeric.empty:
        return False
    return bool(numeric.min() >= 0 and numeric.max() <= 100)


def _style_data_labels(chart, chart_type: str, is_pct: bool):
    """Add smart data labels based on chart type."""
    # Chart types where labels make sense
    label_types = {
        "pie", "doughnut",                                    # always label
        "column_clustered", "bar_clustered",                  # label on top
        "column_stacked", "column_stacked_100",               # label inside
        "bar_stacked", "bar_stacked_100",
        "line", "line_markers",                               # label at points
        "area", "area_stacked",
    }

    if chart_type not in label_types:
        return

    for plot in chart.plots:
        plot.has_data_labels = True
        data_labels = plot.data_labels

        # Number format
        if chart_type in ("column_stacked_100", "bar_stacked_100", "pie", "doughnut"):
            data_labels.number_format = '0.0"%"'
        elif is_pct:
            data_labels.number_format = '0.0"%"'
        else:
            data_labels.number_format = '#,##0.0'

        # Font
        data_labels.font.size = Pt(8)
        data_labels.font.color.rgb = DARK_GRAY

        # Position varies by chart type
        if chart_type in ("pie", "doughnut"):
            data_labels.show_category_name = False
            data_labels.show_percentage = True
            data_labels.show_value = False
            data_labels.font.size = Pt(9)
            data_labels.font.bold = True
        elif "stacked" in chart_type and "100" not in chart_type:
            data_labels.show_value = True
        elif chart_type in ("line", "line_markers"):
            # Only show labels if ≤5 series (otherwise 200+ overlapping labels)
            try:
                n_series = len(chart.plots[0].series) if chart.plots else 0
            except Exception:
                n_series = 0
            if n_series > 5:
                plot.has_data_labels = False
            else:
                data_labels.show_value = True
                data_labels.font.size = Pt(7)
        else:
            data_labels.show_value = True


def _style_series_colors(chart, chart_type: str):
    """Apply professional color palette to chart series."""
    try:
        is_line_type = chart_type in ("line", "line_markers")
        for plot in chart.plots:
            for idx, series in enumerate(plot.series):
                color = CHART_COLORS[idx % len(CHART_COLORS)]
                if is_line_type:
                    # Line charts: ONLY set line color/width — NEVER fill
                    # (fill.solid() on a line fills the area under it → appears opaque/empty)
                    series.format.line.color.rgb = color
                    series.format.line.width = Pt(2.0)
                elif chart_type in ("area", "area_stacked"):
                    # Area charts: fill IS appropriate
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = color
                    series.format.line.color.rgb = color
                    series.format.line.width = Pt(1.5)
                else:
                    # Column/bar/pie: normal fill
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = color
    except Exception:
        pass  # Some chart types don't support series-level coloring


def _style_axes(chart, chart_type: str):
    """Style axis labels and gridlines for readability."""
    try:
        # Category axis (x-axis for column/line, y-axis for bar)
        if hasattr(chart, 'category_axis'):
            cat_ax = chart.category_axis
            cat_ax.tick_labels.font.size = Pt(9)
            cat_ax.tick_labels.font.color.rgb = DARK_GRAY
            cat_ax.has_major_gridlines = False

        # Value axis
        if hasattr(chart, 'value_axis'):
            val_ax = chart.value_axis
            val_ax.tick_labels.font.size = Pt(9)
            val_ax.tick_labels.font.color.rgb = DARK_GRAY
            val_ax.has_major_gridlines = True
            val_ax.major_gridlines.format.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
            val_ax.format.line.color.rgb = RGBColor(0xD9, 0xD9, 0xD9)
    except Exception:
        pass  # Scatter/bubble charts have different axis structure


def render_pptx(df: pd.DataFrame, enc: dict) -> io.BytesIO:
    """Render a native editable PowerPoint chart in-memory.
    
    Returns a BytesIO buffer containing the .pptx file.
    """
    chart_type = enc["chart_type"]
    reg_entry = CHART_REGISTRY[chart_type]

    # Build chart data using the appropriate builder
    builder_name = reg_entry["data_builder"]
    builder_fn = DATA_BUILDERS[builder_name]
    chart_data = builder_fn(df, enc)

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Add a blank slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # ── Slide background (subtle gradient feel via white) ──
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ── Title ──
    title_text = enc.get("title", "Chart")
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT_NAME

    # ── Subtitle (metric context) ──
    if enc.get("rationale"):
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(12), Inches(0.4))
        sub_tf = sub_box.text_frame
        sub_p = sub_tf.paragraphs[0]
        sub_p.text = enc["rationale"]
        sub_p.font.size = Pt(10)
        sub_p.font.italic = True
        sub_p.font.color.rgb = DARK_GRAY
        sub_p.font.name = FONT_NAME

    # ── Chart ──
    x, y, cx, cy = Inches(0.5), Inches(1.2), Inches(12.0), Inches(5.0)
    xl_type = reg_entry["xl_type"]
    chart_shape = slide.shapes.add_chart(xl_type, x, y, cx, cy, chart_data)
    chart = chart_shape.chart

    # ── Legend ──
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    chart.legend.font.color.rgb = DARK_GRAY

    # ── Apply enhancements ──
    is_pct = _is_percentage_data(df, enc)
    _style_data_labels(chart, chart_type, is_pct)
    _style_series_colors(chart, chart_type)
    _style_axes(chart, chart_type)

    # ── Y-axis label ──
    y_label = enc.get("y_axis_label", "")
    if not y_label:
        # Fallback: derive intuitive label from column name
        val_col = enc.get("channels", {}).get("value", "").lower()
        if is_pct or "share" in val_col:
            y_label = "Patient Share (%)" if "patient" in val_col else "Share (%)"
        elif "count" in val_col or "patient" in val_col:
            y_label = "Patient Count"
        elif "volume" in val_col:
            y_label = "Prescription Volume"
        elif "revenue" in val_col or "sales" in val_col:
            y_label = "Revenue ($)"
        elif val_col:
            y_label = val_col.replace("_", " ").title()
    if y_label:
        try:
            chart.value_axis.axis_title.has_text_frame = True
            chart.value_axis.axis_title.text_frame.paragraphs[0].text = y_label
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.name = FONT_NAME
        except Exception:
            pass  # Some chart types (pie, doughnut) have no value axis

    # ── Insight bullets ──
    bullets = enc.get("insight_bullets", [])
    if bullets and any(b.strip() for b in bullets):
        bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.0), Inches(0.7))
        btf = bullet_box.text_frame
        btf.word_wrap = True
        for i, bullet_text in enumerate(bullets[:2]):
            if not bullet_text.strip():
                continue
            if i == 0:
                bp = btf.paragraphs[0]
            else:
                bp = btf.add_paragraph()
            bp.text = f"• {bullet_text.strip()}"
            bp.font.size = Pt(10)
            bp.font.color.rgb = DARK_GRAY
            bp.font.name = FONT_NAME
            bp.space_after = Pt(4)

    # ── Footer ──
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(10), Inches(0.4))
    ft = footer.text_frame
    fp = ft.paragraphs[0]
    fp.text = f"Generated by C3PO Analytics | {time.strftime('%Y-%m-%d %H:%M')} | Source: Databricks SQL"
    fp.font.size = Pt(8)
    fp.font.color.rgb = MID_GRAY
    fp.font.name = FONT_NAME

    # ── Row count badge ──
    badge = slide.shapes.add_textbox(Inches(10.5), Inches(7.0), Inches(2.5), Inches(0.4))
    bt = badge.text_frame
    bp = bt.paragraphs[0]
    bp.text = f"{len(df):,} rows"
    bp.font.size = Pt(8)
    bp.font.color.rgb = MID_GRAY
    bp.font.name = FONT_NAME
    from pptx.enum.text import PP_ALIGN
    bp.alignment = PP_ALIGN.RIGHT

    # Save to in-memory buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    print(f"[PPT] Generated in-memory ({len(buffer.getvalue())} bytes)")
    return buffer
